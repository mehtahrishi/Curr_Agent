import os
import uuid
from flask import Flask, render_template, request, jsonify, session
from werkzeug.utils import secure_filename
from dotenv import load_dotenv
import google.generativeai as genai
from pptx import Presentation
import PyPDF2
import re
import logging
from flask_session import Session
import cloudinary
import cloudinary.uploader
import cloudinary.api
import requests
import io
import redis

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)

# --- Configure Flask-Session ---
app.config["SESSION_TYPE"] = "redis"
app.config["SESSION_PERMANENT"] = False
app.config["SESSION_USE_SIGNER"] = True
app.config["SESSION_REDIS"] = redis.from_url(os.getenv("REDIS_URL"))

# --- Flask Secret Key ---
app.secret_key = os.getenv("FLASK_SECRET_KEY", "hrishi")
if app.secret_key == "hrishi":
    app.logger.warning("Using default FLASK_SECRET_KEY. Please set a strong secret key in your .env file.")

Session(app)

# --- Configure Cloudinary ---
try:
    cloudinary.config(
        cloud_name=os.getenv("CLOUDINARY_CLOUD_NAME"),
        api_key=os.getenv("CLOUDINARY_API_KEY"),
        api_secret=os.getenv("CLOUDINARY_API_SECRET"),
        secure=True
    )
    app.logger.info("Cloudinary configured successfully.")
except Exception as e:
    app.logger.error(f"Error configuring Cloudinary: {e}")

UPLOAD_FOLDER = 'temp_uploads'
ALLOWED_EXTENSIONS = {'pptx', 'pdf'}
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# --- Configure Gemini API ---
try:
    gemini_api_key = os.getenv("GOOGLE_API_KEY")
    if not gemini_api_key:
        raise ValueError("GOOGLE_API_KEY not found in environment variables.")
    genai.configure(api_key=gemini_api_key)
    model = genai.GenerativeModel('gemini-1.5-flash-latest')
except Exception as e:
    if hasattr(app, 'logger'): app.logger.error(f"Error configuring Gemini API: {e}")
    else: print(f"CRITICAL Error configuring Gemini API before logger setup: {e}")
    model = None

def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def extract_text_from_pptx(file_stream):
    try:
        prs = Presentation(file_stream)
        slides_data = []
        for i, slide in enumerate(prs.slides):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slide_text += run.text + " "
            slides_data.append({"type": "slide", "number": i + 1, "content": slide_text.strip()})
        return slides_data
    except Exception as e:
        app.logger.error(f"Error extracting text from PPTX stream: {e}", exc_info=True)
        return None

def extract_text_from_pdf(file_stream):
    pdf_data = []
    try:
        reader = PyPDF2.PdfReader(file_stream)
        if reader.is_encrypted:
            try:
                reader.decrypt('')
            except Exception as decrypt_error:
                app.logger.warning(f"Could not decrypt PDF stream: {decrypt_error}")
                return [{"type": "page", "number": 1, "content": "Error: PDF is encrypted and could not be decrypted."}]
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text is None:
                page_text = "[Content not extractable from this page]"
            pdf_data.append({"type": "page", "number": i + 1, "content": page_text})
    except Exception as e:
        app.logger.error(f"Error extracting text from PDF stream: {e}", exc_info=True)
        return None
    return pdf_data

def get_specific_slide_or_page_content(document_data, user_question):
    patterns = [r"(slide|page)\s*(?:number)?\s*(\d+)", r"(\d+)(?:st|nd|rd|th)?\s*(slide|page)"]
    for pattern in patterns:
        match = re.search(pattern, user_question.lower())
        if match:
            groups = match.groups()
            item_type = groups[1] if groups[0].isdigit() else groups[0]
            item_number = int(groups[0]) if groups[0].isdigit() else int(groups[1])
            for item in document_data:
                if item.get("type") == item_type and item.get("number") == item_number:
                    return item.get("content")
    return None

@app.route('/')
def index():
    current_doc_filename = session.get('original_filename', None)
    return render_template('index.html', current_doc_filename=current_doc_filename)

@app.route('/clear_session', methods=['POST'])
def clear_session_route():
    session.pop('document_data', None)
    session.pop('original_filename', None)
    return jsonify({'status': 'success', 'message': 'Document context cleared.'})

@app.route('/upload', methods=['POST'])
def upload_file_route():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    uploaded_file = request.files['file']
    if uploaded_file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    if uploaded_file and allowed_file(uploaded_file.filename):
        original_filename = secure_filename(uploaded_file.filename)
        try:
            unique_public_id = f"curr_ai_uploads/{uuid.uuid4()}_{original_filename}"
            upload_result = cloudinary.uploader.upload(
                uploaded_file.stream,
                public_id=unique_public_id,
                resource_type="raw",
                folder="curr_ai_uploads"
            )
            cloudinary_url = upload_result.get('secure_url')
            if not cloudinary_url:
                return jsonify({'error': 'File upload to cloud storage failed.'}), 500
            response = requests.get(cloudinary_url, stream=True)
            response.raise_for_status()
            file_content_stream = io.BytesIO(response.content)
            extracted_data = None
            if original_filename.lower().endswith('.pptx'):
                extracted_data = extract_text_from_pptx(file_content_stream)
            elif original_filename.lower().endswith('.pdf'):
                extracted_data = extract_text_from_pdf(file_content_stream)
            file_content_stream.close()
            if not extracted_data or all(not item.get("content", "").strip() for item in extracted_data):
                return jsonify({'error': 'Could not parse document or document appears to be empty.'}), 500
            session['document_data'] = extracted_data
            session['original_filename'] = original_filename
            return jsonify({'status': 'success', 'filename': original_filename, 'message': 'File processed successfully.'})
        except cloudinary.exceptions.Error as e_cld:
            return jsonify({'error': f'Cloudinary API Error: {str(e_cld)}'}), 500
        except requests.exceptions.RequestException as e_req:
            return jsonify({'error': f'Error downloading file from cloud: {str(e_req)}'}), 500
        except Exception as e:
            return jsonify({'error': f'Server error during processing: {str(e)}'}), 500
    else:
        return jsonify({'error': 'Invalid file type. Only PPTX and PDF are allowed.'}), 400

@app.route('/chat', methods=['POST'])
def chat_route():
    if not model:
        return jsonify({'answer': 'AI Chatbot is not configured. Please check API key.'}), 500
    data = request.get_json()
    user_question = data.get('question', '').strip()
    if 'document_data' not in session:
        return jsonify({'answer': 'It seems the document context was lost. Please upload your document again.'}), 400
    if not user_question:
        return jsonify({'answer': 'Please type a question.'}), 400
    document_data = session['document_data']
    doc_filename = session.get('original_filename', 'the current document')
    context_for_llm = ""
    specific_content = get_specific_slide_or_page_content(document_data, user_question)
    MAX_CONTEXT_CHARS = 15000
    if specific_content:
        trimmed_specific_content = specific_content[:MAX_CONTEXT_CHARS]
        if len(specific_content) > MAX_CONTEXT_CHARS: trimmed_specific_content += " [Content Truncated]"
        context_for_llm = f"The user is asking about a specific part of the document '{doc_filename}'. Here is the relevant text from that section:\n---\n{trimmed_specific_content}\n---"
        instruction = f"You are Curr Ai Agent. Based *only* on the provided text excerpt from '{doc_filename}', answer the user's question. If the answer cannot be found in this specific excerpt, clearly state that the information is not in this particular section."
    else:
        full_doc_text_parts = []
        current_chars = 0
        for item in document_data:
            item_header = f"{item.get('type','Item').capitalize()} {item.get('number','N/A')}:\n"
            item_content = item.get('content', '').strip()
            if current_chars + len(item_header) + len(item_content) < MAX_CONTEXT_CHARS:
                full_doc_text_parts.append(item_header + item_content + "\n---\n")
                current_chars += len(item_header) + len(item_content) + 5
            else:
                remaining_space = MAX_CONTEXT_CHARS - current_chars - len(item_header)
                if remaining_space > 50: full_doc_text_parts.append(item_header + item_content[:remaining_space-20] + " [Content Truncated]\n---\n")
                break
        context_for_llm = f"The user is asking a general question about the document '{doc_filename}'. Here is an overview or relevant parts of the document:\n---\n{''.join(full_doc_text_parts)}\n---"
        instruction = f"You are Curr Ai Agent. Based *only* on the provided document context from '{doc_filename}', answer the user's question. If the information is not in the provided context, clearly state that."
    prompt = f"{instruction}\n\nDOCUMENT CONTEXT:\n{context_for_llm}\n\nUSER QUESTION:\n{user_question}\n\nYOUR ANSWER:"
    try:
        chat_session = model.start_chat(history=[])
        response = chat_session.send_message(prompt)
        ai_answer = response.text
    except Exception as e:
        ai_answer = f"Sorry, I encountered an error trying to process your request with the AI model: {str(e)}"
        if "API key not valid" in str(e) or "PERMISSION_DENIED" in str(e):
            ai_answer = "Sorry, there's an issue with the AI service configuration (e.g., API key or permissions). Please contact the administrator."
        elif "rate limit" in str(e).lower():
            ai_answer = "The AI service is currently busy due to high demand. Please try again in a few moments."
    return jsonify({'answer': ai_answer})

if __name__ == '__main__':
    app.logger.setLevel(logging.DEBUG)
    app.run(debug=True, host='0.0.0.0', port=5000)